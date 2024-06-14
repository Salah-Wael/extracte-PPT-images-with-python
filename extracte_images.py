import os
from pptx import Presentation

# Correct relative path to the PowerPoint file
file_name = 'Python_ Introduction.pptx'
folder_path = 'D:\\Computer Science\\my-github\\Python\\extracte-PPT-images-with-python'
file_path = os.path.abspath(folder_path + '\\' + file_name)

# Print current working directory
print(f"Current working directory: {os.getcwd()}")

# Print the contents of the directory
directory = os.path.dirname(file_path)
print(f"Contents of the directory '{directory}':")
try:
    print(os.listdir(directory))
except FileNotFoundError as e:
    print(e)

# Check if the file exists
if not os.path.isfile(file_path):
    raise FileNotFoundError(f"The file {file_path} does not exist.")

# Load the presentation
presentation = Presentation(file_path)

# Create a directory to save images
images_dir = folder_path + '\\' +'extracted_images from ' + file_name[0:-5] + 'PPT'
os.makedirs(images_dir, exist_ok=True)

# Function to save images from the presentation
def save_images(slide, slide_number):
    image_number = 1
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Shape type 13 is a picture
            image = shape.image
            image_bytes = image.blob
            image_filename = os.path.join(images_dir, f'slide{slide_number}_image{image_number}.{image.ext}')
            with open(image_filename, 'wb') as img_file:
                img_file.write(image_bytes)
            image_number += 1

# Iterate through slides and fetch images
for i, slide in enumerate(presentation.slides, start=1):
    print(f"Slide {i}:")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(shape.text)
    save_images(slide, i)

print(f"Images have been saved in the '{images_dir}' directory.")
# to run => python extracte_images.py